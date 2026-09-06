"""生成《山海行 — 文旅多Agent行程规划系统》完整答辩/交付 PPT（v3 完整版）。

覆盖全部内容：项目背景 / 痛点 / 目标需求 / 技术难点 / 亮点 / 架构 /
8 Agent / 核心技术 / 数据模型 / API 前端 / 部署 / 流程演示 / 测试验收 /
设计决策 / 个性化人群适配 / 总结。

排版要点：
- 16:9，统一页头 + 页脚页码
- 中文微软雅黑字体（保证 PowerPoint 里中文不跑版）
- 内容超过 6 条时自动两列，避免文字溢出
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配色
PRIMARY = RGBColor(0x16, 0x77, 0xFF)
PRIMARY_DARK = RGBColor(0x0A, 0x3D, 0x91)
DARK = RGBColor(0x1F, 0x2D, 0x3D)
GREEN = RGBColor(0x52, 0xC4, 0x1A)
RED = RGBColor(0xF5, 0x22, 0x2D)
ORANGE = RGBColor(0xFA, 0xAD, 0x14)
GRAY = RGBColor(0x8C, 0x8C, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF5, 0xFF)
CARD_BG = RGBColor(0xF7, 0xFA, 0xFC)
BORDER = RGBColor(0xE5, 0xEA, 0xF0)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)
FONT = "微软雅黑"

TOTAL = 21


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = BORDER
    shp.shadow.inherit = False
    return shp


def _set_run_font(run, size, color, bold):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    # 显式设置东亚字体，保证中文用微软雅黑
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_text(slide, left, top, width, height, text, size=14, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(3)
        p.line_spacing = line_spacing
        for r in p.runs:
            _set_run_font(r, size, color, bold)
    return box


def header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.32), Inches(11), Inches(0.6), title, size=25, color=DARK, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.92), Inches(11), Inches(0.4), subtitle, size=12, color=GRAY)
    add_rect(slide, Inches(0.5), Inches(1.30), Inches(2.2), Inches(0.045), PRIMARY)


def footer(slide, idx):
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.4), f"{idx}/{TOTAL}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.4), "山海行 · 文旅多 Agent 行程规划系统", size=10, color=GRAY)


def cover(prs, img_path):
    slide = blank(prs)
    set_bg(slide, PRIMARY_DARK)
    if img_path:
        try:
            slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)
        except Exception:
            pass
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x0A, 0x16, 0x28))
    add_text(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2), "山海行 · 文旅多 Agent 行程规划系统", size=42, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.8), "基于 8 个 AI Agent 协作的文旅资源调研与个性化行程规划", size=20, color=RGBColor(0xD0, 0xE4, 0xFF))
    add_text(slide, Inches(1.0), Inches(4.25), Inches(11.3), Inches(0.6), "真实数据 · 断点续传 · 图记忆 · 运行态强干预 · HITL 人工审核 · 人群适配 · RAG 检索增强", size=15, color=RGBColor(0xB0, 0xD0, 0xF0))
    add_text(slide, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5), "版本 v3.2 · 2026-09-05 · 1 人独立完成", size=13, color=RGBColor(0x90, 0xB0, 0xD0))


def toc(prs):
    slide = blank(prs)
    set_bg(slide, WHITE)
    header(slide, "目录")
    items = [
        "01  项目背景", "02  项目痛点", "03  项目目标与需求", "04  技术难点",
        "05  项目亮点", "06  系统架构", "07  8 个 Agent 详解", "08  核心技术实现",
        "09  数据源与数据模型", "10  API 与前端", "11  部署与启动", "12  完整流程演示",
        "13  测试与验收", "14  设计决策", "15  个性化人群适配", "16  RAG 检索增强生成",
        "17  总结",
    ]
    half = 9
    for i, it in enumerate(items):
        col = 0 if i < half else 1
        row = i if i < half else i - half
        left = Inches(1.0 + col * 6.1)
        top = Inches(1.7 + row * 0.62)
        add_text(slide, left, top, Inches(5.6), Inches(0.5), it, size=15, color=DARK)
    footer(slide, 2)


def bullets_slide(prs, title, items, idx, subtitle=""):
    """通用内容页：标题 + 卡片式 bullet，超过 6 条自动两列。"""
    slide = blank(prs)
    set_bg(slide, WHITE)
    header(slide, title, subtitle)
    n = len(items)
    two_col = n > 6
    if two_col:
        left_n = (n + 1) // 2
        cols = [items[:left_n], items[left_n:]]
        col_w = Inches(6.1)
        gap_x = Inches(6.35)
        top0 = Inches(1.55)
        for c, col_items in enumerate(cols):
            m = len(col_items)
            card_h = Inches((5.35 - (m - 1) * 0.08) / m)
            for i, item in enumerate(col_items):
                ct = top0 + i * (card_h + Inches(0.08))
                lx = Inches(0.5) + c * gap_x
                add_rect(slide, lx, ct, col_w, card_h, CARD_BG)
                add_rect(slide, lx, ct, Inches(0.08), card_h, PRIMARY)
                if isinstance(item, tuple):
                    t, content = item
                    add_text(slide, lx + Inches(0.2), ct + Inches(0.1), col_w - Inches(0.4), Inches(0.35), t, size=13, color=PRIMARY_DARK, bold=True)
                    add_text(slide, lx + Inches(0.2), ct + Inches(0.45), col_w - Inches(0.4), card_h - Inches(0.55), content, size=10.5, color=DARK, line_spacing=1.0)
                else:
                    add_text(slide, lx + Inches(0.2), ct + Inches(0.12), col_w - Inches(0.4), card_h - Inches(0.24), item, size=12, color=DARK, line_spacing=1.0)
    else:
        top0 = Inches(1.55)
        gap = Inches(0.08)
        card_h = Inches((5.35 - (n - 1) * 0.08) / n)
        for i, item in enumerate(items):
            ct = top0 + i * (card_h + gap)
            add_rect(slide, Inches(0.5), ct, Inches(12.3), card_h, CARD_BG)
            add_rect(slide, Inches(0.5), ct, Inches(0.08), card_h, PRIMARY)
            if isinstance(item, tuple):
                t, content = item
                add_text(slide, Inches(0.85), ct + Inches(0.1), Inches(11.8), Inches(0.4), t, size=14, color=PRIMARY_DARK, bold=True)
                add_text(slide, Inches(0.85), ct + Inches(0.5), Inches(11.8), Inches(card_h - Inches(0.55)), content, size=12, color=DARK, line_spacing=1.05)
            else:
                add_text(slide, Inches(0.85), ct + Inches(0.15), Inches(11.8), Inches(card_h - Inches(0.3)), item, size=13, color=DARK, line_spacing=1.05)
    footer(slide, idx)


def main():
    import os
    base = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    img = os.path.join(base, "frontend", "src", "assets", "login-bg.png")
    out = os.path.join(base, "docs", "项目讲解.pptx")

    prs = new_prs()

    # 封面
    cover(prs, img)
    # 目录
    toc(prs)

    # 一、项目背景
    bullets_slide(prs, "一、项目背景", [
        ("业务背景", "传统文旅靠人工查资料 + Excel 排行程，一个行程从需求到方案需几小时到一整天：手动搜景点、查口碑、排行程、算预算、查风险。"),
        ("三大业务痛点", "① 数据静态过时（查到的信息可能已过期）② 无风险拦截（可能推荐差评景点/超预算/夜行路段）③ 无容灾（查到一半崩溃全丢）。"),
        ("技术背景", "人工智能课程「多 Agent 协同」工单任务，分两期：多Agent-2（文旅规划主体）+ 工单7（图记忆与运行态强干预叠加）。"),
        ("工期约束", "1 人独立完成，v1.0 共 7 人日 + 工单7 增量 3 人日，需做大量「最小增量、最大复用」的取舍决策。"),
    ], 3)

    # 二、项目痛点
    bullets_slide(prs, "二、项目痛点", [
        ("痛点1：语义遗忘", "用户第 1 次说「海鲜严重过敏」被记住，两周后再规划，系统忘了又推荐海鲜 —— 缺长期记忆（工单7 点名痛点）。"),
        ("痛点2：状态不刷新", "运行中某景点「台风停运」，系统不知道仍按旧信息编排 —— 缺运行态纠偏能力。"),
        ("痛点3：长任务无容灾", "8 个 Agent 跑几分钟，中途进程崩溃，前面工作全丢 —— 缺断点续传。"),
        ("痛点4：无风险拦截", "人工排行程可能推荐差评景点、超预算 45%、高危夜行路段 —— 缺 HITL 人机门控。"),
        ("痛点5：假数据假智能", "早期用写死种子数据 + Mock LLM，景点路线推理全是假的 —— 缺真实数据接入。"),
        ("痛点6：无人群适配", "不管老人还是儿童，行程千篇一律：老人被推荐去滑雪、儿童去酒吧、早餐排火锅、门票不按年龄分档。"),
        ("痛点7：知识不沉淀", "景区政策/地方美食靠写死代码或 LLM 记忆，改起来要改代码、会编造会过时 —— 缺 RAG 检索增强。"),
    ], 4)

    # 三、项目目标与需求
    bullets_slide(prs, "三、项目目标与需求", [
        ("一句话目标", "把「人工排行程」改造成「AI 自动调研工厂」：输入一句自然语言，8 个 AI Agent 自动完成调研、编排、预算、报告。"),
        ("核心功能", "① 8 Agent 协同编排 ② planning-with-files 断点续传 ③ 五状态机+原子写入 ④ HITL 人工审核 ⑤ 安全（内容过滤+防注入）。"),
        ("工单7叠加", "① 共享图记忆（跨行程沉淀知识图谱，修复语义遗忘）② 运行态强干预（主管带验签修改运行中状态）。"),
        ("人群个性化", "老人/儿童门票按年龄分档 + 交通购票按年龄分档 + 成人关系 + 儿童/老人景点过滤 + 兴趣标签推荐。"),
        ("RAG 检索增强", "pgvector + 百炼 embedding + 文旅知识库（6类语料）+ AI 目的地解读。"),
        ("性能红线", "QPS≥200 / P95<300ms / 错误率<0.1% / 检索<150ms / 干预成功率=100%。"),
    ], 5)

    # 四、技术难点（上：前4个）
    bullets_slide(prs, "四、技术难点（1/2）", [
        ("断点续传正确性", "难点：进程随时可能崩溃(kill -9)，怎么保证不丢进度、不重复执行？解法：原子写入(temp→fsync→os.replace→目录fsync)保证任何时刻文件要么旧版要么新版，绝不半份；checksum 校验失败回退快照；已落库任务幂等不重跑。"),
        ("异步队列崩溃接管", "难点：Worker 崩了正在处理的消息怎么办？解法：Redis Streams consumer group + PEL(未确认消息清单) + XAUTOCLAIM(自动认领)，崩溃 worker 的未确认消息自动转给存活 worker。"),
        ("审核并发防重", "难点：两个主管同时审批同一行程，怎么防重复？解法：Lua CAS 抢占(状态原子转换) + DB 条件更新兜底 + 会话租约锁(带 token 防误删)，锁降冲突、DB 兜底保正确。"),
        ("写入-入队竞态", "难点：API 写入行程后立即入队，Worker 抢在 DB commit 前消费，报「行程不存在」。解法：入队前显式 await db.commit()。本地单进程永远发现不了，只有容器化多进程才暴露。"),
    ], 6)

    # 四、技术难点（下：后3个）
    bullets_slide(prs, "四、技术难点（2/2）", [
        ("图记忆抽取检索", "难点：怎么把「我香菜过敏」抽成结构化三元组并跨行程召回？解法：真实 LLM 抽取(不限过敏原白名单) + 三元组 upsert + 向量/图双路检索 + Redis 缓存(0.3ms)。"),
        ("强干预安全原子性", "难点：主管改运行中状态，既不越权又不并发错乱？解法：HMAC 验签 + nonce 防重放 + SETNX 排他锁 + 三写一事务 + 节点入口应用补丁，干预成功率 100%。"),
        ("真实数据接入的坑", "难点：第三方 API 各种坑(key 类型/参数格式/限流/编码)。解法(实践踩坑)：高德 key 需拼参数、百炼 json_object 需含 json 字、nginx 缓存旧 IP 致 502 需 DNS 动态解析、百炼 embedding batch 上限是 10 不是 16。"),
    ], 7)

    # 五、项目亮点
    bullets_slide(prs, "五、项目亮点", [
        ("全链路真实数据", "高德(景点/路线/天气/营业时间/评分/图片) + 百炼LLM(qwen3.7-flash)，无假数据假智能。"),
        ("人群个性化适配", "老人按年龄分档门票 + 成人关系(情侣/家庭等) + 儿童/老人景点过滤 + 兴趣标签推荐 + 餐饮三餐规范(早餐不吃正餐、三餐不重复)。"),
        ("完整行程规划", "每日三餐餐饮 + 景点图片(可放大) + 营业时间 + 门票/人均价 + 真实完成时间 + 出发/返程日期自动算天数。"),
        ("性能大幅超标", "QPS 1109(5.5x)、P95 35ms(快8x)、检索0.3ms(快500x)、干预成功率100%，单行程约18秒。"),
        ("工程细节扎实", "原子写+checksum+快照、Lua CAS+DB兜底、防抖+幂等+超时+心跳、密钥安全注入。"),
        ("工单7全量落地", "图记忆+强干预，数据模型/引擎/API/前端/测试/文档完整闭环。"),
        ("多角色+多语言+可编辑", "游客/顾问/主管三角色 RBAC + 11国语言切换 + 手动编辑/删除行程。"),
        ("精致前端视觉", "真实风景壁纸 + 状态看板友好化(进度条+步骤卡片) + 品牌名山海行。"),
    ], 8)

    # 六、系统架构
    bullets_slide(prs, "六、系统架构（6 层）", [
        ("前端层", "React 18 + TypeScript + AntD，6 个页面（登录/工作台/行程列表/详情/审核台/记忆图谱），WebSocket 实时推送。"),
        ("API 网关层", "FastAPI + JWT/RBAC + 统一错误码 + HMAC 验签，只做入队/查询，不执行长任务。"),
        ("业务服务层", "高德服务、原子写、分布式锁、队列、恢复、安全过滤等可复用逻辑。"),
        ("Agent 编排层", "LangGraph StateGraph + 8 节点 DAG + 条件边 + 记忆抽取 + 干预补丁读取。"),
        ("数据层", "PostgreSQL(pgvector,11表权威) + Redis(Streams/锁/缓存) + 文件(断点恢复源) + 知识库向量。"),
        ("外部数据源", "高德地图(地理编码/POI/路径/天气/图片) + 百炼 LLM(qwen3.7-flash) + Embedding(text-embedding-v3)。"),
    ], 9)

    # 七、8个Agent详解（上：前4个）
    bullets_slide(prs, "七、8 个 Agent 详解（1/2）", [
        ("① Intake 偏好解析", "输入：用户需求(表单/语音/对话)。输出：结构化偏好(目的地/天数/预算/出行人明细/兴趣/成人关系)。技术：真实 LLM + JSON Schema 结构化输出，失败降级正则。附带记忆抽取沉淀过敏/偏好。"),
        ("② Planner 任务拆解", "输入：偏好。输出：≥10 条调研任务(景点/公园/博物馆/酒店/餐厅/早餐/快餐/购物/交通)。技术：CoT 思维链拆解，内置调研维度降级，幂等(已拆过不重复)。"),
        ("③ Web Research 调研", "输入：任务列表。输出：高德真实 POI(名称/地址/评分/营业时间/图片) + 知识库检索结果。技术：ReAct 逐项调研 + 安全过滤 + 断点续传 + RAG 检索注入。"),
        ("④ Sentiment 舆情评估", "输入：真实 POI 列表。输出：风险标签(差评/强制购物/夜行等)。技术：规则引擎关键词打分，识别严重负面舆情供 HITL 触发。"),
    ], 10)

    # 七、8个Agent详解（下：后4个）
    bullets_slide(prs, "七、8 个 Agent 详解（2/2）", [
        ("⑤ Itinerary 日程编排", "输入：POI + 偏好 + 出发地。输出：每日行程(上午/下午/晚上 + 三餐) + 往返交通 + 路线。技术：高德路径规划 + 人群过滤 + 门票分档 + 兴趣排序 + 读取干预补丁。"),
        ("⑥ Budget 预算计算", "输入：行程 + 天数。输出：分项预算(交通/住宿/门票/餐饮/其他) + 超支比例。技术：老人/儿童按年龄分档门票 + 虚拟房价，超预算 45% 触发审核。"),
        ("⑦ Human Review 审核", "输入：风险上下文(超预算/夜行/舆情)。输出：挂起人工审批。技术：不调 LLM，只做路由；主管抢占(Lua CAS)→通过续跑/驳回终止；审批超时 24h 自动驳回。"),
        ("⑧ Report 报告生成", "输入：全部前序输出。输出：完整报告(概览/天气/交通/每日行程/餐饮/酒店/预算/数据来源)。技术：Markdown 渲染 + RAG 生成「目的地深度解读」+ 真实完成时间。"),
    ], 11)

    # 八、核心技术实现
    bullets_slide(prs, "八、核心技术实现", [
        ("断点续传", "temp→fsync→os.replace→目录fsync，任何时刻崩溃文件要么旧版要么新版，绝无半份；checksum 校验失败回退快照。"),
        ("五状态机", "planning→running→suspended→completed + recovering/failed/cancelled，终态禁止再跑图。"),
        ("HITL 门控", "超预算45%/高危夜行/严重舆情 → 挂起 → 主管抢占(Lua CAS) → 通过续跑/驳回终止；审批超时24h自动驳回。"),
        ("图记忆", "LLM抽取三元组→upsert→向量+图双路检索→Redis缓存(0.3ms)，支持任意过敏原(香菜/花粉/青霉素)。"),
        ("强干预", "HMAC验签→nonce防重放→SETNX锁→三写一事务→节点入口应用补丁，干预成功率100%。"),
        ("城市规范化", "高德 geocode 把「河北沧州」规范为「沧州」，修复 POI 搜索误搜到省一级。"),
        ("安全", "JWT+RBAC(游客/顾问/主管) + 注入检测 + 内容过滤 + 工具白名单。"),
    ], 12)

    # 九、数据源与数据模型
    bullets_slide(prs, "九、数据源与数据模型", [
        ("高德地图(1 key 6能力)", "地理编码/POI搜索/路径规划/天气/营业时间/景点图片，全真实数据。"),
        ("阿里云百炼 LLM", "qwen3.7-flash，出发地/目的地解析 + 交通票价估算 + 记忆抽取 + RAG 生成（Mock 可切真实）。"),
        ("百炼 Embedding", "text-embedding-v3（768 维，OpenAI 兼容），语义向量化（图记忆 + 知识库）。"),
        ("虚拟房价", "真实酒店 POI 名 + 城市等级分档估算价（一线/二线/其他三档）。"),
        ("虚拟门票/人均", "景点门票按年龄分档（老人60-64半价/65+免首道、儿童6岁以下免/6-18半价）+ 餐厅人均参考价。"),
        ("业务表(6张)", "users / travel_plans / agent_tasks / review_records / budget_records / audit_logs。"),
        ("图记忆表(4张)", "graph_nodes / graph_edges / memory_events / interventions。"),
        ("知识库表(1张)", "document_chunks（6类语料，38 chunk，Vector 768）。"),
    ], 13)

    # 十、API与前端
    bullets_slide(prs, "十、API 与前端", [
        ("API(5组20+接口)", "认证(register含角色/login) / 行程(create/list/detail/status/agents/plan-file/report/cancel/delete/PATCH itinerary) / 审核 / 记忆(工单7) / 运维。"),
        ("记忆接口(工单7)", "intervene(强干预) / rollback(回滚) / graph(子图) / search(检索) / interventions(历史)。"),
        ("前端 8 页面", "登录注册 / 工作台 / 行程列表 / 行程详情 / 审核台 / 记忆图谱 / 语音创建 / 对话式创建。"),
        ("行程详情亮点", "出发地交通方式+票价 + 调研结果(真实POI) + 景点图片(可放大) + 三餐餐饮 + 门票按年龄分档 + 优待备注 + 编辑行程。"),
        ("创建表单亮点", "成人关系 + 老人信息(年龄+性别+状态) + 儿童信息(年龄+身高+占座) + 学生学历 + 购票/酒店方式。"),
        ("多种创建方式", "语音创建(浏览器原生语音识别) + 对话式创建(AI 主动引导，像 DeepSeek 一样聊天)。"),
    ], 14)

    # 十一、部署与启动
    bullets_slide(prs, "十一、部署与启动", [
        ("一键启动", "双击 一键启动.bat → 自动构建+启动+打开浏览器；一键关闭.bat 停止；数据保留在 Docker 卷。"),
        ("7 个 Docker 服务", "api / worker / frontend / postgres / redis / prometheus / grafana。"),
        ("访问地址", "前端 8080 / API 文档 8001 / Grafana 3002。"),
        ("演示账号", "advisor_demo(顾问) / supervisor_demo(主管)，密码 wenlv123；注册可自选游客/顾问身份。"),
        ("多国语言", "左下角切换 11 种语言（中/英/日/韩/法/德/西/俄/葡/意/阿）。"),
        ("nginx 根治", "Docker DNS 动态解析，api 容器重建不 502。"),
    ], 15)

    # 十二、完整流程演示
    bullets_slide(prs, "十二、完整流程演示（5 场景）", [
        ("场景A 正常生成", "提交「3天杭州游」→ 8 Agent 执行(真实LLM+高德) → 报告含天气+路线+行程+餐饮+酒店+预算。"),
        ("场景B 崩溃恢复", "Web Research 跑到第6页 kill -9 → 重启 → 从第7页续跑，前6页不重复请求。"),
        ("场景C 注入拦截", "第7页含「忽略之前指令」对抗样本 → 拦截 → 标记 blocked → 写审计日志。"),
        ("场景D 人工审核", "超预算/夜行/舆情 → 挂起 → 主管抢占 → 通过 → 续跑出报告。"),
        ("场景E 图记忆+干预", "行程1说「海鲜过敏」→ 行程2自动避开；主管干预「台风停运」→ 运行中自动剔除。"),
        ("场景F 人群适配", "带65岁老人+5岁儿童 → 老人门票免首道、儿童免票；剔除滑雪/酒吧；早餐吃地方特色早餐。"),
        ("场景G RAG增强", "检索知识库(门票政策+云南美食) → 报告生成「目的地深度解读」，有据可查、带人群感知。"),
    ], 16)

    # 十三、测试与验收
    bullets_slide(prs, "十三、测试与验收", [
        ("压测结果", "QPS 1109(≥200 ✅ 5.5x) / P95 35.4ms(<300 ✅ 快8x) / 错误率 0%(<0.1% ✅)。"),
        ("工单7 红线", "检索 P95 0.3ms(<150 ✅ 快500x) / 干预成功率 100%(=100% ✅)。"),
        ("单元测试", "8 个用例全绿(test_memory_intervention.py：三元组抽取/验签/embedding)。"),
        ("联调脚本", "e2e_check / recovery_test / simulate_conversation(15轮) / concurrent(50并发) / verify_retrieve_latency。"),
        ("竞态审计", "50 并发干预零丢失更新(Lost Update=0)，版本链完整可回放。"),
        ("RAG 检索验证", "老人门票政策 0.83 / 早餐吃什么 0.69 / 儿童高铁买票 0.83 / 海鲜过敏 0.67。"),
    ], 17)

    # 十四、设计决策
    bullets_slide(prs, "十四、设计决策（关键取舍）", [
        ("编排框架", "LangGraph（工单首选，7人日不允许自研图引擎）。"),
        ("队列", "Redis Streams（PEL 原生断点，比 Celery 轻）。"),
        ("LLM", "抽象 Provider，Mock→真实切换（压测不烧钱，演示切真模型）。"),
        ("图存储", "PostgreSQL + pgvector（不引 Neo4j，单机原则）。"),
        ("Embedding", "百炼 text-embedding-v3（768 维，复用现有 key，OpenAI 兼容）。"),
        ("记忆引擎", "自研轻量 Mem0 风格（3人日不允许消化外部框架黑盒）。"),
        ("强干预", "HMAC + nonce + 三写事务（工单红线：安全验签）。"),
    ], 18)

    # 十五、个性化人群适配
    bullets_slide(prs, "十五、个性化人群适配", [
        ("老人门票分档", "按年龄：60 以下全价 / 60-64 半价 / 65+ 免首道大门票（国有A级景区规则），只免首道门票，观光车/索道/演出另计。"),
        ("儿童门票分档", "6周岁以下或身高1.2米以下免首道大门票（二选一）、6-18半价；只免首道，民营景区以公示为准。"),
        ("交通购票分档", "高铁(6岁以下免票不占座/6-14半价/14+全价)；飞机(婴儿票10%/儿童5折/12+成人)；学生票仅限家校往返。"),
        ("人群信息表单", "成人关系 + 老人(年龄+性别+状态) + 儿童(年龄+身高+占座) + 学生(学历) + 购票/酒店方式。"),
        ("人群过滤+兴趣", "有老人剔除滑雪/登山/雪山；有儿童剔除酒吧/夜店；兴趣标签+亲子/老人友好加权排序。"),
        ("餐饮三餐规范", "早餐只吃早餐类(不吃火锅/烤鸭)；早餐/快餐专门调研；三餐跨天不重复。"),
    ], 19)

    # 十六、RAG 检索增强生成
    bullets_slide(prs, "十六、RAG 检索增强生成", [
        ("技术底座", "pgvector(PostgreSQL扩展) + 百炼 text-embedding-v3(768维) + OpenAI 兼容 /embeddings，batch=10 规避单次上限。"),
        ("场景A 个人记忆", "图记忆 embedding 从 JSONB 换 Vector(768)，MD5哈希假向量换真语义向量，「海鲜过敏」能召回「对虾过敏」。"),
        ("场景B 文旅知识库", "document_chunks 表存 6 类语料(景点攻略/地方美食/景区政策/四季玩法/交通贴士/住宿贴士)，38 个分块。"),
        ("检索(R)", "search_kb() 用 pgvector 余弦距离召回 top-k 原文 chunk；web_research + itinerary 两节点注入。"),
        ("生成(G)", "report 节点 _generate_rag_insights() 让 LLM 基于原文生成「目的地深度解读」，有据可查、带人群感知。"),
        ("降级安全", "非 real 模式或检索为空或 LLM 失败时返回 None，报告照常出（只少一节），不阻塞主流程。"),
    ], 20)

    # 十七、总结
    slide = blank(prs)
    set_bg(slide, LIGHT)
    add_rect(slide, 0, Inches(2.1), Inches(0.25), Inches(1.6), PRIMARY)
    add_text(slide, Inches(0.9), Inches(2.1), Inches(11), Inches(1.0), "总结", size=30, color=DARK, bold=True)
    pts = [
        ("完成度", "原始需求(多Agent-2) + 工单7(图记忆+强干预) + 人群适配 + RAG 检索增强 全部落地，验收级完整。"),
        ("数据真实性", "全链路真实数据（高德+百炼 LLM+Embedding），无假数据假智能。"),
        ("个性化能力", "老人/儿童门票按年龄分档 + 交通购票分档 + 人群过滤 + 兴趣推荐 + 餐饮规范。"),
        ("RAG 能力", "pgvector + 语义向量 + 文旅知识库 + AI 目的地解读，检索增强生成完整闭环。"),
        ("性能", "四项指标全部达标且大幅超标（QPS 5.5x / 检索快500x）。"),
        ("核心能力", "断点续传 + HITL + 图记忆 + 强干预 + 安全 + 人群适配 + RAG，七维完整。"),
    ]
    top = Inches(2.9)
    for i, (t, c) in enumerate(pts):
        ct = top + i * Inches(0.68)
        add_rect(slide, Inches(0.9), ct, Inches(11.5), Inches(0.6), WHITE)
        add_rect(slide, Inches(0.9), ct, Inches(0.08), Inches(0.6), PRIMARY)
        add_text(slide, Inches(1.2), ct + Inches(0.06), Inches(2.0), Inches(0.5), t, size=13, color=PRIMARY_DARK, bold=True)
        add_text(slide, Inches(3.2), ct + Inches(0.06), Inches(9.0), Inches(0.5), c, size=11, color=DARK)
    footer(slide, 21)

    prs.save(out)
    print(f"PPT 已生成：{out}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
